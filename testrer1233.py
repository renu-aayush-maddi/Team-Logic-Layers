"""
json_to_ppt.py

Creates a PPTX from JSON like:
{
  "title": "Profile of M Renu Aayush",
  "slides": [
    {
      "slide_title": "Introduction",
      "content": [
        "Name: M Renu Aayush",
        "Email: renuaayushm@gmail.com",
        ...
      ]
    },
    ...
  ]
}
"""

import json
import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

SAMPLE_JSON = {
  "title": "Profile of M Renu Aayush",
  "slides": [
    {
      "slide_title": "Introduction",
      "content": [
        "Name: M Renu Aayush",
        "Email: renuaayushm@gmail.com",
        "Contact: +91-9396763455",
        "Education: Undergraduate at KL University",
        "Field: Computer Science and Engineering"
      ]
    },
    {
      "slide_title": "Academic Background",
      "content": [
        "Under the guidance of Professor Sagar Imambi",
        "Roll Number: 2200032294",
        "Department: CSE (DST-FIST Sponsored Department)",
        "Koneru Lakshmaiah Education Foundation",
        "Green Fields, Vaddeswaram, Guntur District - 522 502"
      ]
    },
    {
      "slide_title": "Project Certificate",
      "content": [
        "Project Title: Metal Surface Defect Detection",
        "Course: 22AIP3305/A, Deep Learning",
        "Certification of bonafide work and fulfillment of Bachelor's Degree requirements",
        "Project claimed as original work with no submission to other universities"
      ]
    },
    {
      "slide_title": "Achievements",
      "content": [
        "Flipkart Grid 5.0 Finalist - Selected from 400,000+ participants",
        "BOSCH x Tinkerers’ Lab Hackathon – IIT Hyderabad – Finalist (Top 15 out of 100 teams)",
        "Odyssey of Code – ConsultAdd Hackathon – Secured 6th place; recognized for innovation"
      ]
    },
    {
      "slide_title": "Acknowledgments",
      "content": [
        "Thanks to Professor Sagar Imambi for outstanding support",
        "Gratitude to Dr P. Vidyullatha, Course Co-ordinator for 22AIP3305R/A/P Deep Learning"
      ]
    },
    {
      "slide_title": "Projects",
      "content": [
        "Developed collaborative archival systems for preserving Indian cultural artifacts",
        "Implemented role-based access control with Spring Security",
        "Created a fully functional e-commerce website using the MERN stack (Bharat Coin Bazaar)",
        "Improved user experience by 25% through personalized interactions"
      ]
    }
  ]
}

def create_ppt_from_json(data: dict, output_path: str = "output_presentation.pptx"):
    prs = Presentation()

    # Title slide layout (usually layout 0)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    # Title text
    title_text = data.get("title", "Presentation")
    title.text = title_text

    # Optional subtitle (blank or could add generation info)
    if subtitle is not None:
        subtitle.text = ""

    # Content slides
    slides = data.get("slides", [])
    # Prefer layout 1 or 5 depending on template — use a Title and Content layout
    content_layout = None
    for i in (1, 5, 6, 3):  # try common indices
        try:
            content_layout = prs.slide_layouts[i]
            break
        except IndexError:
            continue
    if content_layout is None:
        content_layout = prs.slide_layouts[1]

    for s in slides:
        slide_title = s.get("slide_title", "")
        content_items = s.get("content", [])

        slide = prs.slides.add_slide(content_layout)
        # set the title
        try:
            slide.shapes.title.text = slide_title
        except Exception:
            pass

        # find the first placeholder that accepts text (content placeholder)
        body = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                body = shape
                break
        # fallback: create a textbox
        if body is None:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            body = txBox

        tf = body.text_frame
        tf.clear()  # clear any existing text

        # Add each content item as a bullet paragraph
        for i, item in enumerate(content_items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = str(item)
            p.level = 0  # top-level bullet
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            # font sizing and properties
            for run in p.runs:
                run.font.size = Pt(18)
                # run.font.name = 'Calibri'  # optional

        # If there were no content items, keep a small hint
        if not content_items:
            p = tf.paragraphs[0]
            p.text = "(no content)"
            for run in p.runs:
                run.font.size = Pt(14)

    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")

def load_json_from_file(path: str):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def main():
    # Usage:
    # python json_to_ppt.py           -> uses embedded SAMPLE_JSON
    # python json_to_ppt.py data.json -> uses data.json
    # python json_to_ppt.py data.json out.pptx -> custom output filename
    args = sys.argv[1:]
    data = SAMPLE_JSON
    output = "output_presentation.pptx"
    if len(args) >= 1:
        try:
            data = load_json_from_file(args[0])
        except Exception as e:
            print(f"Failed to load {args[0]}: {e}\nFalling back to sample JSON.")
            data = SAMPLE_JSON
    if len(args) >= 2:
        output = args[1]

    create_ppt_from_json(data, output)

if __name__ == "__main__":
    main()
