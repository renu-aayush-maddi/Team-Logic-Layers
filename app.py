<<<<<<< HEAD
# import os
# import asyncio
# import json
# import tempfile
# from typing import List, Optional
# from fastapi import FastAPI, Request, HTTPException, UploadFile, File, Form
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.staticfiles import StaticFiles
# from dotenv import load_dotenv

# from langchain_openai import ChatOpenAI
# from src.helper import load_pdf_file, text_split
# from src.prompt import structured_system_prompt
# from langchain_core.prompts import ChatPromptTemplate
# from langchain_openai import OpenAIEmbeddings
# from langchain_chroma import Chroma
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from langchain.chains import create_retrieval_chain

# load_dotenv()

# app = FastAPI()
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"]
# )
# app.mount("/static", StaticFiles(directory="static"), name="static")

# OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
# if not OPENAI_API_KEY:
#     raise RuntimeError("OPENAI_API_KEY is required in environment")

# llm = ChatOpenAI(openai_api_key=OPENAI_API_KEY, model="gpt-4o-mini")
# prompt_chroma = ChatPromptTemplate.from_messages([
#     ("system", structured_system_prompt),
#     ("human", "{input}")
# ])

# @app.post("/api/v1/hackrx/run")
# async def run_query(
#     request: Request,
#     file: UploadFile = File(...),
#     questions: Optional[str] = Form(None)  # either JSON array string, single question string, or omitted if multiple question fields used
# ):
#     """
#     file: uploaded PDF
#     questions: either
#       - JSON array string: '["Q1","Q2"]'
#       - single question string: 'What is X?'
#       - omitted, and multiple form entries named 'questions' supplied (request.form().getlist('questions'))
#     """
#     # --- parse questions robustly ---
#     try:
#         q_list = None

#         # 1) if `questions` single field provided, try parse it as JSON array
#         if questions is not None:
#             # try JSON array
#             try:
#                 parsed = json.loads(questions)
#                 if isinstance(parsed, list) and all(isinstance(x, str) for x in parsed):
#                     q_list = parsed
#                 else:
#                     # not an array -> treat whole string as one question
#                     q_list = [str(questions)]
#             except json.JSONDecodeError:
#                 # not JSON -> treat as single question
#                 q_list = [str(questions)]
#         else:
#             # 2) maybe multiple form fields with same name: questions=Q1&questions=Q2
#             form = await request.form()
#             # form.getlist exists and returns list of values for repeated keys
#             try:
#                 repeated = form.getlist("questions")
#             except Exception:
#                 # fallback: collect all values named 'questions'
#                 repeated = [v for k, v in form.multi_items() if k == "questions"]
#             if repeated:
#                 q_list = [str(x) for x in repeated]

#         if not q_list:
#             raise HTTPException(status_code=400, detail="No questions provided. Send 'questions' as JSON array string or multiple 'questions' form fields.")
#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

#     # --- save uploaded PDF to a temp file ---
#     temp_filepath = None
#     try:
#         suffix = os.path.splitext(file.filename)[1] or ".pdf"
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
#             content = await file.read()
#             if not content:
#                 raise HTTPException(status_code=400, detail="Uploaded file is empty")
#             tf.write(content)
#             temp_filepath = tf.name

#         # --- extract pages (list[Document]) ---
#         extracted_docs = load_pdf_file(temp_filepath)
#         if not isinstance(extracted_docs, list) or len(extracted_docs) == 0:
#             # no extracted text — give helpful error
#             raise HTTPException(status_code=422, detail="No text extracted from PDF. If your PDF is a scanned image, run OCR (e.g., tesseract) to create a searchable PDF before upload.")

#         # --- chunking: helper.text_split(extracted_docs, pdf_path) ---
#         # Ensure helper supports this signature. If it expects raw text, adapt accordingly.
#         chunks = text_split(extracted_docs, pdf_path=temp_filepath)

#         # --- embeddings + Chroma ---
#         embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY)
#         chroma_dir = "./chroma_db"
#         os.makedirs(chroma_dir, exist_ok=True)
#         docsearch = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_dir)

#         retriever = docsearch.as_retriever(search_type="similarity", search_kwargs={"k": 15})
#         qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
#         rag_chain = create_retrieval_chain(retriever, qa_chain)

#         async def process_question(q: str):
#             resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
#             if isinstance(resp, dict):
#                 return resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
#             return str(resp)

#         answers = await asyncio.gather(*[process_question(q) for q in q_list])

#         return {"answers": answers}

#     except HTTPException:
#         # re-raise HTTP errors
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
#     finally:
#         # cleanup tempfile
#         try:
#             if temp_filepath and os.path.exists(temp_filepath):
#                 os.remove(temp_filepath)
#         except Exception:
#             pass


import asyncio
import json
import easyocr
import os
import tempfile
from typing import List, Optional

from dotenv import load_dotenv
from fastapi import (
    FastAPI,
    HTTPException,
    Request,
    UploadFile,
    File,
    Form,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles

from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_chroma import Chroma
from langchain_openai import ChatOpenAI, OpenAIEmbeddings

from src.helper import load_pdf_file, text_split
from src.prompt import structured_system_prompt

# add imports near top of file
import os
import tempfile
import asyncio
from typing import Optional, List
from fastapi import Request, UploadFile, File, Form, HTTPException

# OCR/image libs
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes


load_dotenv()

reader = easyocr.Reader(["en"])
=======
import os
import asyncio
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv

from langchain_google_genai import GoogleGenerativeAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_pinecone import PineconeVectorStore
from pinecone import Pinecone

from src.prompt import structured_system_prompt
from src.helper import download_hugging_face_embeddings

load_dotenv()

>>>>>>> 809cc7dbcad3a17c96f6084921cee17d01c24b55
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
<<<<<<< HEAD
    allow_headers=["*"],
)
app.mount("/static", StaticFiles(directory="static"), name="static")

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY is required in environment")

# LLM and prompt setup
llm = ChatOpenAI(openai_api_key=OPENAI_API_KEY, model="gpt-4o-mini")
prompt_chroma = ChatPromptTemplate.from_messages(
    [("system", structured_system_prompt), ("human", "{input}")]
)


def _parse_questions_from_form(questions_field: Optional[str], request: Request) -> List[str]:
    """Robustly parse questions coming from the form.

    Accepts:
    - JSON array string: '["Q1", "Q2"]'
    - Single question string: 'What is X?'
    - Omitted, but multiple form fields with same name 'questions' supplied
    """
    # 1) If questions single field provided, try parse it as JSON array
    if questions_field is not None:
        try:
            parsed = json.loads(questions_field)
            if isinstance(parsed, list) and all(isinstance(x, str) for x in parsed):
                return parsed
            # not an array -> treat whole string as one question
            return [str(questions_field)]
        except json.JSONDecodeError:
            # not JSON -> treat as single question
            return [str(questions_field)]

    # 2) Maybe multiple form fields with same name
    # Note: this helper is synchronous, so we must access the request form in the caller (async)
    raise ValueError("No questions in questions_field; caller must fetch repeated form fields")


def ocr_image_bytes(image_bytes: bytes, tesseract_lang: str = "eng") -> str:
    try:
        with Image.open(tempfile.SpooledTemporaryFile()) as _:
            pass
    except Exception:
        # fallback: open from bytes directly
        pass
    try:
        img = Image.open(io.BytesIO(image_bytes))
    except Exception as e:
        # some images may be multipart, try temp file
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".img")
        tmp.write(image_bytes)
        tmp.close()
        img = Image.open(tmp.name)

    # convert to RGB to avoid mode issues
    if img.mode != "RGB":
        img = img.convert("RGB")
    text = pytesseract.image_to_string(img, lang=tesseract_lang)
    return text

# --- helper: OCR a PDF with scanned pages -> return list of page texts ---
def ocr_pdf_bytes(pdf_bytes: bytes, dpi: int = 300, tesseract_lang: str = "eng") -> List[str]:
    # convert PDF bytes to images (one per page)
    try:
        pages = convert_from_bytes(pdf_bytes, dpi=dpi)
    except Exception as e:
        raise RuntimeError(f"pdf->image conversion failed: {str(e)}")
    page_texts = []
    for page in pages:
        # ensure RGB
        if page.mode != "RGB":
            page = page.convert("RGB")
        text = pytesseract.image_to_string(page, lang=tesseract_lang)
        page_texts.append(text)
    return page_texts

# make sure to import io
import io

# from fastapi.responses import JSONResponse

# @app.post("/api/v1/hackrx/run")
# async def run_query(
#     request: Request,
#     file: UploadFile = File(...),
#     questions: Optional[str] = Form(None),
# ):
#     """Handle uploaded PDF and answer one or more questions using a RAG flow.
#     Returns JSONResponse(content={"ok": True, "parsed": parsed, "raw": llm_output, "files_payload_summary": summary})
#     """

#     # --- parse questions robustly ---
#     try:
#         try:
#             q_list = _parse_questions_from_form(questions, request)
#         except ValueError:
#             form = await request.form()
#             try:
#                 repeated = form.getlist("questions")
#             except Exception:
#                 repeated = [v for k, v in form.multi_items() if k == "questions"]
#             if repeated:
#                 q_list = [str(x) for x in repeated]
#             else:
#                 raise HTTPException(
#                     status_code=400,
#                     detail=(
#                         "No questions provided. Send 'questions' as JSON array string or "
#                         "multiple 'questions' form fields."
#                     ),
#                 )
#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

#     temp_filepath: Optional[str] = None

#     try:
#         suffix = os.path.splitext(file.filename)[1] or ".pdf"
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
#             content = await file.read()
#             if not content:
#                 raise HTTPException(status_code=400, detail="Uploaded file is empty")
#             tf.write(content)
#             temp_filepath = tf.name

#         # --- extract pages (list[Document]) ---
#         extracted_docs = load_pdf_file(temp_filepath)
#         if not isinstance(extracted_docs, list) or len(extracted_docs) == 0:
#             raise HTTPException(
#                 status_code=422,
#                 detail=(
#                     "No text extracted from PDF. If your PDF is a scanned image, run OCR "
#                     "(e.g., tesseract) to create a searchable PDF before upload."
#                 ),
#             )

#         # --- chunking ---
#         chunks = text_split(extracted_docs, pdf_path=temp_filepath)

#         # --- embeddings + Chroma ---
#         embeddings = OpenAIEmbeddings(
#             model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY
#         )
#         chroma_dir = "./chroma_db"
#         os.makedirs(chroma_dir, exist_ok=True)

#         docsearch = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_dir)
#         retriever = docsearch.as_retriever(
#             search_type="similarity", search_kwargs={"k": 15}
#         )

#         qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
#         rag_chain = create_retrieval_chain(retriever, qa_chain)

#         # process_question now returns a tuple: (parsed_answer, raw_response)
#         async def process_question(q: str):
#             resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
#             # prefer common keys but keep raw response as-is
#             if isinstance(resp, dict):
#                 parsed_ans = resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
#             else:
#                 parsed_ans = str(resp)
#             return parsed_ans, resp

#         results = await asyncio.gather(*[process_question(q) for q in q_list])

#         # Split parsed answers and raw outputs into separate lists
#         parsed = [r[0] for r in results]           # cleaned strings for callers
#         llm_output = [r[1] for r in results]       # raw LLM/RAG responses (dict or str)

#         # Lightweight file payload summary (adjust fields as you like)
#         summary = {
#             "filename": file.filename,
#             "temp_filepath": temp_filepath,
#             "num_extracted_docs": len(extracted_docs),
#             "num_chunks": len(chunks),
#         }

#         return JSONResponse(
#             content={
#                 "ok": True,
#                 "parsed": parsed,
#                 "raw": llm_output,
#                 "files_payload_summary": summary,
#             }
#         )

#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
#     finally:
#         try:
#             if temp_filepath and os.path.exists(temp_filepath):
#                 os.remove(temp_filepath)
#         except Exception:
#             pass
"""
app.py

FastAPI backend that:
- Accepts JSON via POST /convert
- auto-detects whether JSON is for PPTX (slides) or PDF (sections)
- or honors query param ?format=pptx or ?format=pdf
- returns a downloadable PPTX or PDF bytes (StreamingResponse)

Dependencies:
pip install fastapi uvicorn python-multipart python-pptx reportlab
"""

import io
import json
from typing import Any, Dict, Optional

from fastapi import FastAPI, Request, HTTPException, Query
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware

# pptx/reportlab imports
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

app = FastAPI(title="JSON → PPTX/PDF Converter")

# Allow CORS from frontend dev server (adjust origins as needed)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # in prod narrow this to your frontend origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def is_ppt_json(data: Dict[str, Any]) -> bool:
    # Heuristic: "slides" array -> PPT
    return isinstance(data.get("slides"), list)


def is_pdf_json(data: Dict[str, Any]) -> bool:
    # Heuristic: "sections" array OR having a title but not slides -> PDF-style doc
    if isinstance(data.get("sections"), list):
        return True
    if "slides" not in data and "title" in data:
        return True
    return False


def simple_wrap_text(text: str, max_chars: int = 100):
    words = text.split()
    lines = []
    cur = ""
    for w in words:
        if len(cur) + len(w) + 1 <= max_chars:
            cur = (cur + " " + w).strip()
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def create_pptx_from_json(data: Dict[str, Any]) -> bytes:
    prs = Presentation()

    # Title slide
    title_text = data.get("title", "Presentation")
    try:
        title_layout = prs.slide_layouts[0]
    except Exception:
        title_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    try:
        slide.shapes.title.text = str(title_text)
    except Exception:
        pass
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = ""
        except Exception:
            pass

    slides = data.get("slides", [])
    # pick a standard content layout if available
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    for s in slides:
        slide_title = s.get("slide_title", "")
        # content may be list or string; support both
        content_items = s.get("content") or s.get("body") or s.get("content_items") or []

        cur_slide = prs.slides.add_slide(content_layout)
        try:
            cur_slide.shapes.title.text = str(slide_title)
        except Exception:
            pass

        # find first placeholder with text frame
        tf = None
        for shape in cur_slide.shapes:
            if shape.is_placeholder and getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                break

        if tf is None:
            left = Inches(0.5)
            top = Inches(1.6)
            width = Inches(9)
            height = Inches(5)
            txBox = cur_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame

        tf.clear()

        if isinstance(content_items, list):
            for i, item in enumerate(content_items):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = str(item)
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(14)
        else:
            p = tf.paragraphs[0]
            p.text = str(content_items)
            for run in p.runs:
                run.font.size = Pt(14)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()


def create_pdf_from_json(data: Dict[str, Any]) -> bytes:
    bio = io.BytesIO()
    c = canvas.Canvas(bio, pagesize=A4)
    width, height = A4

    margin_x = 50
    current_y = height - 50

    # Title
    title = data.get("title")
    if title:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(margin_x, current_y, str(title))
        current_y -= 30

    sections = data.get("sections", [])
    if isinstance(sections, list) and sections:
        for sec in sections:
            sec_title = sec.get("title") if isinstance(sec, dict) else None
            sec_body = sec.get("body") if isinstance(sec, dict) else (sec if isinstance(sec, str) else None)

            if sec_title:
                c.setFont("Helvetica-Bold", 14)
                if current_y < 80:
                    c.showPage()
                    current_y = height - 50
                c.drawString(margin_x, current_y, str(sec_title))
                current_y -= 20

            if sec_body:
                c.setFont("Helvetica", 12)
                lines = simple_wrap_text(str(sec_body), max_chars=100)
                for ln in lines:
                    if current_y < 60:
                        c.showPage()
                        current_y = height - 50
                    c.drawString(margin_x, current_y, ln)
                    current_y -= 16
            current_y -= 10
            if current_y < 80:
                c.showPage()
                current_y = height - 50
    else:
        # fallback: if slides present, render each slide as a PDF page
        slides = data.get("slides", [])
        if slides:
            for s in slides:
                slide_title = s.get("slide_title", "")
                content_items = s.get("content", []) or s.get("body", [])
                c.setFont("Helvetica-Bold", 16)
                c.drawString(margin_x, current_y, str(slide_title))
                y = current_y - 30
                c.setFont("Helvetica", 12)
                for item in (content_items or []):
                    lines = simple_wrap_text(str(item), max_chars=100)
                    for ln in lines:
                        if y < 60:
                            c.showPage()
                            y = height - 50
                        c.drawString(margin_x, y, ln)
                        y -= 16
                c.showPage()
        else:
            # dump raw JSON pretty
            c.setFont("Helvetica", 10)
            pretty = json.dumps(data, indent=2)
            y = current_y
            for line in pretty.splitlines():
                if y < 60:
                    c.showPage()
                    y = height - 50
                c.drawString(margin_x, y, line[:110])
                y -= 12

    c.save()
    bio.seek(0)
    return bio.read()


@app.post("/convert")
async def convert(request: Request, format: Optional[str] = Query("auto")):
    """
    POST /convert?format=auto|pdf|pptx
    Body: raw JSON
    Returns: file (pptx or pdf)
    """
    try:
        data = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON payload")

    # If explicit format requested, respect it
    fmt = (format or "auto").lower()
    try:
        if fmt == "pptx":
            content = create_pptx_from_json(data)
            return StreamingResponse(io.BytesIO(content),
                                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'})
        elif fmt == "pdf":
            content = create_pdf_from_json(data)
            return StreamingResponse(io.BytesIO(content),
                                    media_type="application/pdf",
                                    headers={"Content-Disposition": 'attachment; filename="document.pdf"'})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Generation error: {e}")

    # auto detect
    if is_ppt_json(data):
        content = create_pptx_from_json(data)
        return StreamingResponse(io.BytesIO(content),
                                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'})
    else:
        # default PDF
        content = create_pdf_from_json(data)
        return StreamingResponse(io.BytesIO(content),
                                media_type="application/pdf",
                                headers={"Content-Disposition": 'attachment; filename="document.pdf"'})


# optional healthcheck
@app.get("/health")
async def health():
    return {"status": "ok"}

# @app.post("/api/v1/hackrx/run")
# async def run_query(
#     request: Request,
#     file: UploadFile = File(...),
#     questions: Optional[str] = Form(None),
# ):
#     """Handle an uploaded PDF and answer one or more questions using a RAG flow.

#     Parameters
#     ----------
#     file: UploadFile
#         Uploaded PDF file (searchable PDF expected). If the PDF is scanned images only,
#         the handler will return a helpful error asking the user to run OCR first.

#     questions: Optional[str]
#         Either a JSON array string ("[\"Q1\", \"Q2\"]"), a single question string,
#         or omitted (in which case multiple form fields named 'questions' are expected).
#     """

#     # --- parse questions robustly ---
#     try:
#         try:
#             q_list = _parse_questions_from_form(questions, request)
#         except ValueError:
#             form = await request.form()
#             # form.getlist exists and returns list of values for repeated keys
#             try:
#                 repeated = form.getlist("questions")
#             except Exception:
#                 # fallback: collect all values named 'questions'
#                 repeated = [v for k, v in form.multi_items() if k == "questions"]
#             if repeated:
#                 q_list = [str(x) for x in repeated]
#             else:
#                 raise HTTPException(
#                     status_code=400,
#                     detail=(
#                         "No questions provided. Send 'questions' as JSON array string or "
#                         "multiple 'questions' form fields."
#                     ),
#                 )
#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

#     # --- save uploaded PDF to a temp file ---
#     temp_filepath: Optional[str] = None

#     try:
#         suffix = os.path.splitext(file.filename)[1] or ".pdf"
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
#             content = await file.read()
#             if not content:
#                 raise HTTPException(status_code=400, detail="Uploaded file is empty")
#             tf.write(content)
#             temp_filepath = tf.name

#         # --- extract pages (list[Document]) ---
#         extracted_docs = load_pdf_file(temp_filepath)
#         if not isinstance(extracted_docs, list) or len(extracted_docs) == 0:
#             # no extracted text — give helpful error
#             raise HTTPException(
#                 status_code=422,
#                 detail=(
#                     "No text extracted from PDF. If your PDF is a scanned image, run OCR "
#                     "(e.g., tesseract) to create a searchable PDF before upload."
#                 ),
#             )

#         # --- chunking: helper.text_split(extracted_docs, pdf_path) ---
#         # Ensure helper supports this signature. If it expects raw text, adapt accordingly.
#         chunks = text_split(extracted_docs, pdf_path=temp_filepath)

#         # --- embeddings + Chroma ---
#         embeddings = OpenAIEmbeddings(
#             model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY
#         )
#         chroma_dir = "./chroma_db"
#         os.makedirs(chroma_dir, exist_ok=True)

#         docsearch = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_dir)
#         retriever = docsearch.as_retriever(
#             search_type="similarity", search_kwargs={"k": 15}
#         )

#         qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
#         rag_chain = create_retrieval_chain(retriever, qa_chain)

#         async def process_question(q: str) -> str:
#             # run the (synchronous) chain in a thread to avoid blocking the event loop
#             resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
#             if isinstance(resp, dict):
#                 return resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
#             return str(resp)

#         answers = await asyncio.gather(*[process_question(q) for q in q_list])
#         return {"answers": answers}

#     except HTTPException:
#         # re-raise HTTP errors
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
#     finally:
#         # cleanup tempfile
#         try:
#             if temp_filepath and os.path.exists(temp_filepath):
#                 os.remove(temp_filepath)
#         except Exception:
#             pass


# imports
import os
import io
from typing import List
from google.cloud import vision
from google.oauth2 import service_account
from pdf2image import convert_from_bytes
from PIL import Image

# --- optional: create client from env JSON content OR use default ADC via GOOGLE_APPLICATION_CREDENTIALS ---
def _get_vision_client():
    """
    Preferable: if GOOGLE_APPLICATION_CREDENTIALS is set, google client will use it automatically.
    If you instead stored the JSON content in MY_VISION_SA_JSON, we create credentials from it.
    """
    # If user provided JSON content in MY_VISION_SA_JSON, use it
    json_text = os.environ.get("MY_VISION_SA_JSON")
    if json_text:
        try:
            import json
            info = json.loads(json_text)
            creds = service_account.Credentials.from_service_account_info(info)
            return vision.ImageAnnotatorClient(credentials=creds)
        except Exception as e:
            raise RuntimeError(f"Failed to create Vision client from MY_VISION_SA_JSON: {e}")
    # Otherwise rely on Application Default Credentials (GOOGLE_APPLICATION_CREDENTIALS or cloud runtime)
    return vision.ImageAnnotatorClient()

# --- OCR an image (bytes) via Google Vision text_detection ---
def google_ocr_image_bytes(image_bytes: bytes, max_pages: int = 10) -> str:
    """
    Returns plain text string detected in the image bytes using Google Vision.
    """
    if not image_bytes:
        return ""
    client = _get_vision_client()
    image = vision.Image(content=image_bytes)
    # For single images, text_detection is fine
    response = client.text_detection(image=image)
    if response.error and response.error.message:
        raise RuntimeError(f"Google Vision error: {response.error.message}")
    # response.text_annotations[0] contains the full text
    texts = response.text_annotations
    if not texts:
        return ""
    return texts[0].description or ""

# --- OCR a PDF bytes (scanned pages) by converting pages to images and calling Vision ---
def google_ocr_pdf_bytes(pdf_bytes: bytes, dpi: int = 300, poppler_path: str | None = None) -> List[str]:
    """
    Converts PDF -> images (one per page) using pdf2image, then runs Google OCR per page.
    Returns list of page texts.
    """
    pages = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path=poppler_path)
    client = _get_vision_client()
    page_texts = []
    for page_img in pages:
        try:
            # Convert PIL Image to bytes (JPEG) to pass to Vision
            bio = io.BytesIO()
            page_img_rgb = page_img.convert("RGB")
            page_img_rgb.save(bio, format="JPEG")
            img_bytes = bio.getvalue()
            image = vision.Image(content=img_bytes)
            response = client.document_text_detection(image=image)  # document_text_detection often gives better structured text
            if response.error and response.error.message:
                raise RuntimeError(f"Google Vision PDF page error: {response.error.message}")
            # full text is in response.full_text_annotation.text
            text = ""
            if response.full_text_annotation and response.full_text_annotation.text:
                text = response.full_text_annotation.text
            else:
                # fallback to text_annotations
                if response.text_annotations:
                    text = response.text_annotations[0].description
            page_texts.append(text or "")
        finally:
            try:
                page_img.close()
            except Exception:
                pass
    return page_texts

from google.cloud import vision
import io

def google_ocr_image_bytes(image_bytes: bytes) -> str:
    client = vision.ImageAnnotatorClient()
    image = vision.Image(content=image_bytes)

    response = client.text_detection(image=image)
    if response.error.message:
        raise RuntimeError(f"Google Vision OCR error: {response.error.message}")

    texts = response.text_annotations
    if not texts:
        return ""
    # First element is the full detected text
    return texts[0].description





import tempfile
from pydub import AudioSegment
import openai
try:
    from faster_whisper import WhisperModel
except Exception:
    WhisperModel = None

# initialize openai if key present
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "") or ""
if OPENAI_API_KEY:
    try:
        openai.api_key = OPENAI_API_KEY
    except Exception:
        pass

# optional local faster-whisper model cache
_FW_MODEL = None
def get_faster_whisper_model(size="small"):
    global _FW_MODEL
    if _FW_MODEL is None and WhisperModel is not None:
        try:
            _FW_MODEL = WhisperModel(size)
        except Exception:
            _FW_MODEL = None
    return _FW_MODEL

# ----------------------------
# Audio transcription helpers
# ----------------------------

import os
import tempfile
import asyncio
from typing import Optional, Dict, Any, List
from openai import OpenAI as OpenAIClient  # if you use OpenAI client already

async def transcribe_audio_bytes(content: bytes, filename: str = "upload") -> Dict[str, Any]:
    """
    Try OpenAI transcription (if OPENAI_API_KEY present), else fallback to local whisper/faster-whisper.
    Returns: {"text": full_text, "segments": [ {"start": float, "end": float, "text": "..."} , ... ]}
    """
    # try OpenAI first if key available
    openai_key = os.environ.get("OPENAI_API_KEY") or os.environ.get("OPENAI_API_KEY_OVERRIDE")
    if openai_key:
        try:
            return await _transcribe_with_openai(content, filename, openai_key)
        except Exception as e:
            # failover to local whisper, but log
            # if you have a logger, use logger.exception(...)
            print("OpenAI transcription failed, falling back to local whisper:", e)

    # fallback to local whisper / faster-whisper
    return await _transcribe_with_local_whisper(content, filename)


async def _transcribe_with_openai(content: bytes, filename: str, openai_key: str) -> Dict[str, Any]:
    """
    Use OpenAI's audio transcription endpoint via the official SDK.
    Synchronous blocking call is run in thread to avoid blocking event loop.
    """
    def sync_call():
        client = OpenAIClient(api_key=openai_key)
        # write to temp file as the SDK expects a file object
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1] or ".mp3") as tf:
            tf.write(content)
            tf.flush()
            tfname = tf.name
        try:
            # model: choose "gpt-4o-mini-transcribe" or "whisper-1" if available in your API plan
            # adjust model name if needed
            resp = client.audio.transcriptions.create(model="gpt-4o-mini-transcribe", file=open(tfname, "rb"))
            # resp structure may be client-specific; commonly resp.text contains transcription
            text = getattr(resp, "text", None) or resp.get("text") if isinstance(resp, dict) else None
            # OpenAI sometimes returns segments; if not, return full text and empty segments
            segments = []
            if hasattr(resp, "segments"):
                segments = resp.segments
            elif isinstance(resp, dict) and "segments" in resp:
                segments = resp["segments"]
            return {"text": text or str(resp), "segments": segments}
        finally:
            try:
                os.remove(tfname)
            except Exception:
                pass

    return await asyncio.to_thread(sync_call)


async def _transcribe_with_local_whisper(content: bytes, filename: str) -> Dict[str, Any]:
    """
    Try faster-whisper first, then openai-whisper.
    Returns dict {"text": ..., "segments": [...]}
    """
    # write to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1] or ".mp3") as tf:
        tf.write(content)
        tf.flush()
        tmp_path = tf.name

    try:
        # First try faster-whisper (recommended if installed)
        try:
            from faster_whisper import WhisperModel
            model_size = os.environ.get("WHISPER_MODEL", "small")  # small|base|tiny|medium|large
            # run in a thread
            def faster_sync():
                # device can be "cpu" or "cuda"
                device = os.environ.get("WHISPER_DEVICE", "cpu")
                model = WhisperModel(model_size, device=device, compute_type="float32")
                segments, info = model.transcribe(tmp_path, beam_size=5, vad_filter=True)
                full_text = ""
                segs = []
                for seg in segments:
                    # segments yield (segment) with start, end, text
                    start = float(seg.start)
                    end = float(seg.end)
                    text = seg.text
                    full_text += text + " "
                    segs.append({"start": start, "end": end, "text": text})
                return {"text": full_text.strip(), "segments": segs}

            return await asyncio.to_thread(faster_sync)

        except Exception as e_fw:
            # fallback to openai-whisper
            print("faster-whisper not available or failed:", e_fw)

        # try openai/whisper package
        try:
            import whisper
            model_name = os.environ.get("WHISPER_MODEL", "small")
            def whisper_sync():
                m = whisper.load_model(model_name)
                result = m.transcribe(tmp_path)
                # result usually contains 'text' and 'segments' if verbose=False/True
                text = result.get("text", "")
                segments_raw = result.get("segments", [])
                segs = []
                for s in segments_raw:
                    segs.append({
                        "start": float(s.get("start", 0.0)),
                        "end": float(s.get("end", 0.0)),
                        "text": s.get("text", "").strip()
                    })
                return {"text": text or "", "segments": segs}
            return await asyncio.to_thread(whisper_sync)
        except Exception as e_w:
            print("openai-whisper not available or failed:", e_w)
            raise RuntimeError("No local whisper backend available. Install 'faster-whisper' or 'openai-whisper', or provide OPENAI_API_KEY for cloud transcription.")
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass


def transcript_to_documents(text: str, segments: Optional[List[Dict[str, Any]]] = None, source_name: str = "audio") -> List[Dict[str, Any]]:
    """
    Convert transcript into a list of documents compatible with text_split.
    Each doc: {"page_content": text_chunk, "metadata": {"source": source_name, "start":..., "end":...}}
    If segments provided, create one doc per segment; else return one doc with entire text.
    """
    docs = []
    if segments:
        for i, seg in enumerate(segments):
            seg_text = seg.get("text", "").strip()
            if not seg_text:
                continue
            meta = {"source": source_name, "segment_index": i}
            if "start" in seg:
                meta["start_seconds"] = seg["start"]
            if "end" in seg:
                meta["end_seconds"] = seg["end"]
            docs.append({"page_content": seg_text, "metadata": meta})
    else:
        docs.append({"page_content": text or "", "metadata": {"source": source_name}})
    return docs


# @app.post("/api/v1/hackrx/run")
# async def run_query(
#     request: Request,
#     file: UploadFile = File(...),
#     questions: Optional[str] = Form(None),
# ):
#     """Handle uploaded PDF, scanned-PDF, or images and answer questions using a RAG flow."""
#     # parse questions (reuse your existing _parse_questions_from_form)
#     try:
#         try:
#             q_list = _parse_questions_from_form(questions, request)
#         except ValueError:
#             form = await request.form()
#             try:
#                 repeated = form.getlist("questions")
#             except Exception:
#                 repeated = [v for k, v in form.multi_items() if k == "questions"]
#             if repeated:
#                 q_list = [str(x) for x in repeated]
#             else:
#                 raise HTTPException(
#                     status_code=400,
#                     detail=(
#                         "No questions provided. Send 'questions' as JSON array string or "
#                         "multiple 'questions' form fields."
#                     ),
#                 )
#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

#     temp_filepath: Optional[str] = None
#     temp_bytes = None
#     try:
#         # read file bytes
#         content = await file.read()
#         if not content:
#             raise HTTPException(status_code=400, detail="Uploaded file is empty")

#         filename = file.filename or "upload"
#         ext = os.path.splitext(filename)[1].lower()

#         # Save a temp file (useful for libraries that need a path)
#         suffix = ext if ext else ".bin"
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
#             tf.write(content)
#             temp_filepath = tf.name

#         extracted_docs = None  # will be list-like expected by your text_split
#         # ---- If PDF: try normal PDF text extraction first ----
#         if ext == ".pdf":
#             # try your existing PDF loader
#             try:
#                 extracted_docs = load_pdf_file(temp_filepath)  # expected list[Document] or similar
#             except Exception as e:
#                 # swallow and let OCR attempt below if no text extracted
#                 extracted_docs = None

#             # if no text extracted (likely scanned PDF), run OCR on PDF pages
#             if not extracted_docs or (isinstance(extracted_docs, list) and len(extracted_docs) == 0):
#                 try:
#                     page_texts = ocr_pdf_bytes(content)
#                 except RuntimeError as e:
#                     raise HTTPException(status_code=422, detail=f"OCR conversion failed: {str(e)}")
#                 # convert page_texts to same shape as load_pdf_file would produce
#                 # minimal Document-like dicts with 'page_content' and 'metadata'
#                 extracted_docs = []
#                 for i, txt in enumerate(page_texts):
#                     if txt and txt.strip():
#                         extracted_docs.append({"page_content": txt, "metadata": {"page": i + 1}})
#                 if len(extracted_docs) == 0:
#                     raise HTTPException(
#                         status_code=422,
#                         detail=(
#                             "No text could be extracted from the PDF even after OCR. "
#                             "Make sure the PDF is legible or try increasing dpi."
#                         ),
#                     )

#         # ---- If image type: run OCR on the image ----
#         elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
#             try:
#                 results = reader.readtext(content, detail=0)  # list of strings
#                 text = "\n".join(results)
#             except Exception as e:
#                 raise HTTPException(status_code=500, detail=f"EasyOCR failed: {e}")
            
#             if not text.strip():
#                 raise HTTPException(
#             status_code=422,
#             detail="No text found in image via EasyOCR."
#                 )
#             extracted_docs = [{"page_content": text, "metadata": {"source": filename}}]


            
        

#         # ---- Other single-file types (optional) ----
#         else:
#             # fallback: try PDF loader anyway (some .pdf disguised)
#             try:
#                 extracted_docs = load_pdf_file(temp_filepath)
#             except Exception:
#                 extracted_docs = None
#             if not extracted_docs:
#                 raise HTTPException(
#                     status_code=415,
#                     detail=(
#                         "Unsupported file type. This endpoint accepts text PDFs, scanned PDFs (images inside PDFs), "
#                         "and image files (.png/.jpg/.jpeg/.tiff). For other files like .docx or audio, use the "
#                         "dedicated endpoints (or ask me to add them)."
#                     ),
#                 )

#         # --- chunking: expect text_split to accept the list of docs and pdf_path if needed ---
#         chunks = text_split(extracted_docs, pdf_path=temp_filepath)

#         # --- embeddings + Chroma (your existing flow) ---
#         embeddings = OpenAIEmbeddings(
#             model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY
#         )
#         chroma_dir = "./chroma_db"
#         os.makedirs(chroma_dir, exist_ok=True)

#         docsearch = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_dir)
#         retriever = docsearch.as_retriever(
#             search_type="similarity", search_kwargs={"k": 15}
#         )

#         qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
#         rag_chain = create_retrieval_chain(retriever, qa_chain)

#         async def process_question(q: str) -> str:
#             resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
#             if isinstance(resp, dict):
#                 return resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
#             return str(resp)

#         answers = await asyncio.gather(*[process_question(q) for q in q_list])
#         return {"answers": answers}

#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
#     finally:
#         try:
#             if temp_filepath and os.path.exists(temp_filepath):
#                 os.remove(temp_filepath)
#         except Exception:
#             pass
        
        
        
from fastapi import Request

# @app.post("/api/v1/hackrx/run")
# async def run_query(
#     request: Request,
#     file: UploadFile = File(...),
#     questions: Optional[str] = Form(None),
# ):
#     """Handle uploaded PDF, scanned-PDF, images, and audio and answer questions using your RAG flow."""
#     # parse questions (same logic you had)
#     try:
#         try:
#             q_list = _parse_questions_from_form(questions, request)
#         except ValueError:
#             form = await request.form()
#             try:
#                 repeated = form.getlist("questions")
#             except Exception:
#                 repeated = [v for k, v in form.multi_items() if k == "questions"]
#             if repeated:
#                 q_list = [str(x) for x in repeated]
#             else:
#                 raise HTTPException(
#                     status_code=400,
#                     detail=(
#                         "No questions provided. Send 'questions' as JSON array string or "
#                         "multiple 'questions' form fields."
#                     ),
#                 )
#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

#     temp_filepath: Optional[str] = None
#     try:
#         content = await file.read()
#         if not content:
#             raise HTTPException(status_code=400, detail="Uploaded file is empty")

#         filename = file.filename or "upload"
#         ext = os.path.splitext(filename)[1].lower()

#         # Write temp file if some loaders need a path (PDF loader uses a path)
#         suffix = ext if ext else ".bin"
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
#             tf.write(content)
#             temp_filepath = tf.name

#         extracted_docs = None

#         # ---- PDF flow (text PDF or scanned PDF via OCR) ----
#         if ext == ".pdf":
#             try:
#                 extracted_docs = load_pdf_file(temp_filepath)  # your existing loader
#             except Exception:
#                 extracted_docs = None

#             if not extracted_docs or (isinstance(extracted_docs, list) and len(extracted_docs) == 0):
#                 # scanned PDF fallback: OCR pages (you already have ocr_pdf_bytes or google_ocr_pdf_bytes)
#                 try:
#                     page_texts = ocr_pdf_bytes(content)
#                 except RuntimeError as e:
#                     raise HTTPException(status_code=422, detail=f"OCR conversion failed: {str(e)}")
#                 extracted_docs = []
#                 for i, txt in enumerate(page_texts):
#                     if txt and txt.strip():
#                         extracted_docs.append({"page_content": txt, "metadata": {"page": i + 1}})
#                 if len(extracted_docs) == 0:
#                     raise HTTPException(status_code=422, detail="No text could be extracted from the PDF even after OCR.")

#         # ---- Image flow (single image file) ----
#         # elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
#         #     try:
#         #         results = reader.readtext(content, detail=0)  # your EasyOCR usage
#         #         text = "\n".join(results)
#         #     except Exception as e:
#         #         raise HTTPException(status_code=500, detail=f"EasyOCR failed: {e}")
#         #     if not text.strip():
#         #         raise HTTPException(status_code=422, detail="No text found in image via EasyOCR.")
#         #     extracted_docs = [{"page_content": text, "metadata": {"source": filename}}]
        
#         # ---- Image flow (single image file) ----
#         elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
#             try:
#         # Path to your service account JSON key (update as needed)
#                 SERVICE_ACCOUNT_JSON = r"C:\Users\renua\OneDrive\Desktop\Template Generation\ocr-rag-project-061fa931dd4a.json"

#         # Load credentials and create Vision client
#                 creds = service_account.Credentials.from_service_account_file(SERVICE_ACCOUNT_JSON)
#                 client = vision.ImageAnnotatorClient(credentials=creds)

#         # Build Vision image from bytes (content is read above)
#                 image = vision.Image(content=content)

#         # Call Vision API for text detection
#                 response = client.text_detection(image=image)
#                 texts = response.text_annotations

#                 if response.error.message:
#             # Vision returned an API error
#                     raise RuntimeError(f"Google Vision API error: {response.error.message}")

#                 if not texts or len(texts) == 0:
#                     raise HTTPException(status_code=422, detail="No text found in image via Google Vision OCR.")

#         # Use the full extracted text (first annotation holds the full text)
#                 text = texts[0].description or ""
#                 print(text)
#             except HTTPException:
#         # re-raise HTTPExceptions to preserve status codes
#                 raise
#             except Exception as e:
#         # map other errors to a 500 so caller knows OCR failed
#                 raise HTTPException(status_code=500, detail=f"GoogleOCR failed: {e}")

#             if not text.strip():
#                 raise HTTPException(status_code=422, detail="No text found in image via Google Vision OCR.")

#     # Produce extracted_docs in the same format as the PDF/OCR flow
#             extracted_docs = [{"page_content": text, "metadata": {"source": filename}}]


#         # ---- Audio flow (merged here) ----
#         elif ext in {".wav", ".mp3", ".m4a", ".flac", ".ogg", ".aac"}:
#             # Transcribe audio (OpenAI Whisper preferred, fallback to faster-whisper)
#             try:
#                 trans = await transcribe_audio_bytes(content, filename=filename)
#             except Exception as e:
#                 raise HTTPException(status_code=500, detail=f"Transcription failed: {e}")
#             text = trans.get("text", "")
#             segments = trans.get("segments", [])
#             if not text or not text.strip():
#                 raise HTTPException(status_code=422, detail="Transcription produced no text")
#             # convert transcript to documents compatible with text_split
#             extracted_docs = transcript_to_documents(text, segments=segments, source_name=filename)

#         # ---- Other single-file types (docx etc.) ----
#         else:
#             try:
#                 extracted_docs = load_pdf_file(temp_filepath)
#             except Exception:
#                 extracted_docs = None
#             if not extracted_docs:
#                 raise HTTPException(
#                     status_code=415,
#                     detail=(
#                         "Unsupported file type. This endpoint accepts text PDFs, scanned PDFs, "
#                         "image files (.png/.jpg/.jpeg/.tiff), and audio files (.wav/.mp3/.m4a/.flac/.ogg/.aac)."
#                     ),
#                 )

#         # --- chunking: pass extracted_docs to your text_split ---
#         chunks = text_split(extracted_docs, pdf_path=temp_filepath)

#         # --- embeddings + Chroma (your existing flow) ---
#         embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY)
#         chroma_dir = "./chroma_db"
#         os.makedirs(chroma_dir, exist_ok=True)

#         docsearch = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_dir)
#         retriever = docsearch.as_retriever(search_type="similarity", search_kwargs={"k": 15})

#         qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
#         rag_chain = create_retrieval_chain(retriever, qa_chain)

#         async def process_question(q: str) -> str:
#             resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
#             if isinstance(resp, dict):
#                 return resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
#             return str(resp)

#         answers = await asyncio.gather(*[process_question(q) for q in q_list])
#         return {"answers": answers}

#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
#     finally:
#         try:
#             if temp_filepath and os.path.exists(temp_filepath):
#                 os.remove(temp_filepath)
#         except Exception:
#             pass
        


import os
import time
import tempfile
import asyncio
from typing import Optional, List, Dict, Any

from fastapi import Request, UploadFile, File, Form, HTTPException

# Google Vision
from google.cloud import vision
from google.oauth2 import service_account

# LangChain / Chroma imports
try:
    from langchain.schema import Document
except Exception:
    Document = None
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import Chroma

DEFAULT_SERVICE_ACCOUNT_JSON = r"C:\Users\renua\OneDrive\Desktop\Template Generation\ocr-rag-project-061fa931dd4a.json"
SERVICE_ACCOUNT_JSON = os.getenv("GOOGLE_APPLICATION_CREDENTIALS", DEFAULT_SERVICE_ACCOUNT_JSON)


@app.post("/api/v1/hackrx/run")
async def run_query(
    request: Request,
    file: UploadFile = File(...),
    questions: Optional[str] = Form(None),
):
    """Handle uploaded PDF, scanned-PDF, images, and audio and answer questions using your RAG flow."""
    # parse questions (same logic you had)
    try:
        try:
            q_list = _parse_questions_from_form(questions, request)
        except ValueError:
            form = await request.form()
            try:
                repeated = form.getlist("questions")
            except Exception:
                repeated = [v for k, v in form.multi_items() if k == "questions"]
            if repeated:
                q_list = [str(x) for x in repeated]
            else:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        "No questions provided. Send 'questions' as JSON array string or "
                        "multiple 'questions' form fields."
                    ),
                )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error parsing questions: {str(e)}")

    temp_filepath: Optional[str] = None
    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Uploaded file is empty")

        filename = file.filename or "upload"
        ext = os.path.splitext(filename)[1].lower()

        # Write temp file if some loaders need a path (PDF loader uses a path)
        suffix = ext if ext else ".bin"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
            tf.write(content)
            temp_filepath = tf.name

        extracted_docs = None

        # ---- PDF flow (text PDF or scanned PDF via OCR) ----
        if ext == ".pdf":
            try:
                extracted_docs = load_pdf_file(temp_filepath)  # your existing loader
            except Exception:
                extracted_docs = None

            if not extracted_docs or (isinstance(extracted_docs, list) and len(extracted_docs) == 0):
                # scanned PDF fallback: OCR pages (you already have ocr_pdf_bytes or google_ocr_pdf_bytes)
                try:
                    page_texts = ocr_pdf_bytes(content)
                except RuntimeError as e:
                    raise HTTPException(status_code=422, detail=f"OCR conversion failed: {str(e)}")
                extracted_docs = []
                for i, txt in enumerate(page_texts):
                    if txt and txt.strip():
                        extracted_docs.append({"page_content": txt, "metadata": {"page": i + 1}})
                if len(extracted_docs) == 0:
                    raise HTTPException(status_code=422, detail="No text could be extracted from the PDF even after OCR.")

        # ---- Image flow (single image file) ----
        elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
            try:
                # Load credentials and create Vision client
                creds = service_account.Credentials.from_service_account_file(DEFAULT_SERVICE_ACCOUNT_JSON)
                client = vision.ImageAnnotatorClient(credentials=creds)

                # Build Vision image from bytes (content is read above)
                image = vision.Image(content=content)

                # Call Vision API for text detection
                response = client.text_detection(image=image)
                texts = getattr(response, "text_annotations", None)

                # robust error check
                if getattr(response, "error", None) and getattr(response.error, "message", None):
                    raise RuntimeError(f"Google Vision API error: {response.error.message}")

                if not texts or len(texts) == 0:
                    raise HTTPException(status_code=422, detail="No text found in image via Google Vision OCR.")

                # Use the full extracted text (first annotation holds the full text)
                text = texts[0].description or ""
                print(f"DEBUG: OCR extracted text length={len(text)}")
            except HTTPException:
                raise
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"GoogleOCR failed: {e}")

            if not text.strip():
                raise HTTPException(status_code=422, detail="No text found in image via Google Vision OCR.")

            # Produce extracted_docs in the same format as the PDF/OCR flow
            extracted_docs = [{"page_content": text, "metadata": {"source": filename}}]

        # ---- Audio flow (merged here) ----
        elif ext in {".wav", ".mp3", ".m4a", ".flac", ".ogg", ".aac"}:
            # Transcribe audio (OpenAI Whisper preferred, fallback to faster-whisper)
            try:
                trans = await transcribe_audio_bytes(content, filename=filename)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Transcription failed: {e}")
            text = trans.get("text", "")
            print(text)
            segments = trans.get("segments", [])
            if not text or not text.strip():
                raise HTTPException(status_code=422, detail="Transcription produced no text")
            # convert transcript to documents compatible with text_split
            extracted_docs = transcript_to_documents(text, segments=segments, source_name=filename)

        # ---- Other single-file types (docx etc.) ----
        else:
            try:
                extracted_docs = load_pdf_file(temp_filepath)
            except Exception:
                extracted_docs = None
            if not extracted_docs:
                raise HTTPException(
                    status_code=415,
                    detail=(
                        "Unsupported file type. This endpoint accepts text PDFs, scanned PDFs, "
                        "image files (.png/.jpg/.jpeg/.tiff), and audio files (.wav/.mp3/.m4a/.flac/.ogg/.aac)."
                    ),
                )

        # --- START: conversion -> chunking -> embeddings -> chroma ---
        # 1) validate
        if not extracted_docs or len(extracted_docs) == 0:
            raise HTTPException(status_code=422, detail="No extracted documents to chunk.")

        # 2) Convert dict-like docs to langchain Documents if needed
        needs_conversion = isinstance(extracted_docs, list) and isinstance(extracted_docs[0], dict) and Document is not None

        docs_for_split = []
        if needs_conversion:
            for d in extracted_docs:
                pc = d.get("page_content", "") if isinstance(d, dict) else ""
                meta = d.get("metadata", {}) if isinstance(d, dict) else {}
                if not str(pc).strip():
                    continue
                docs_for_split.append(Document(page_content=str(pc), metadata=meta))
        else:
            docs_for_split = extracted_docs

        if not docs_for_split or len(docs_for_split) == 0:
            raise HTTPException(status_code=422, detail="After conversion, no non-empty documents found for chunking.")

        # Debug: show input snippets
        print(f"DEBUG: docs_for_split count = {len(docs_for_split)}")
        for i, d in enumerate(docs_for_split[:3]):
            try:
                snippet = d.page_content[:200].replace("\n", " ")
                meta = d.metadata
            except Exception:
                snippet = str(d)[:200].replace("\n", " ")
                meta = getattr(d, "metadata", {}) or {}
            print(f"DEBUG doc {i}: {snippet!r} (meta={meta})")

        # 3) Chunk via your text_split (ensure text_split expects Document or adapt)
        chunks = text_split(docs_for_split, pdf_path=temp_filepath)

        if not chunks or len(chunks) == 0:
            raise HTTPException(status_code=422, detail="Chunking produced zero chunks. Check text_split implementation and input documents.")

        print(f"DEBUG: chunk count = {len(chunks)}")
        for i, c in enumerate(chunks[:3]):
            try:
                content_snip = c.page_content[:200]
                meta = c.metadata
            except Exception:
                content_snip = c.get("page_content", "") if isinstance(c, dict) else str(c)[:200]
                meta = c.get("metadata", {}) if isinstance(c, dict) else {}
            print(f"DEBUG chunk {i}: {str(content_snip).replace(chr(10),' ')}... (meta={meta})")

        # 4) Convert chunks to Documents for Chroma if necessary
        final_docs_for_chroma: List[Any] = []
        if isinstance(chunks[0], dict) and Document is not None:
            for c in chunks:
                pc = c.get("page_content", "")
                meta = c.get("metadata", {})
                if not str(pc).strip():
                    continue
                final_docs_for_chroma.append(Document(page_content=str(pc), metadata=meta))
        else:
            final_docs_for_chroma = chunks

        if not final_docs_for_chroma or len(final_docs_for_chroma) == 0:
            raise HTTPException(status_code=422, detail="No valid chunks to index into Chroma.")

        # 5) Build embeddings and Chroma index
        embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=OPENAI_API_KEY)

        # Use a persistent dir; change while debugging if you want unique DBs
        chroma_dir = "./chroma_db"
        os.makedirs(chroma_dir, exist_ok=True)

        try:
            docsearch = Chroma.from_documents(final_docs_for_chroma, embeddings, persist_directory=chroma_dir)
            # persist if available
            try:
                if hasattr(docsearch, "persist"):
                    docsearch.persist()
            except Exception as e:
                print(f"WARNING: Chroma.persist() raised: {e}")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to create Chroma index: {e}")

        # 6) create retriever / chains
        retriever = docsearch.as_retriever(search_type="similarity", search_kwargs={"k": 15})
        qa_chain = create_stuff_documents_chain(llm, prompt_chroma)
        rag_chain = create_retrieval_chain(retriever, qa_chain)
        # --- END: conversion -> chunking -> embeddings -> chroma ---

        # Ask questions
        async def process_question(q: str) -> str:
            resp = await asyncio.to_thread(rag_chain.invoke, {"input": q})
            if isinstance(resp, dict):
                return resp.get("answer") or resp.get("output") or resp.get("text") or str(resp)
            return str(resp)

        answers = await asyncio.gather(*[process_question(q) for q in q_list])
        return {"answers": answers}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"[RAG ERROR] {str(e)}")
    finally:
        try:
            if temp_filepath and os.path.exists(temp_filepath):
                os.remove(temp_filepath)
        except Exception:
            pass
=======
    allow_headers=["*"]
)

# API Keys
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
os.environ["PINECONE_API_KEY"] = PINECONE_API_KEY

# Pinecone + retriever
embeddings_default = download_hugging_face_embeddings()
index_name = "test"
Pinecone(api_key=PINECONE_API_KEY)
docsearch_default = PineconeVectorStore.from_existing_index(index_name=index_name, embedding=embeddings_default)
retriever_default = docsearch_default.as_retriever(search_type="similarity", search_kwargs={"k": 3})

# LLM setup
llm_default = GoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.4, max_tokens=500)
prompt_default = ChatPromptTemplate.from_messages([
    ("system", structured_system_prompt),
    ("human", "{input}\n\nAdditional Info (if any): {location_info}")
])
qa_chain_default = create_stuff_documents_chain(llm_default, prompt_default)
rag_chain_default = create_retrieval_chain(retriever_default, qa_chain_default)

class QuestionRequest(BaseModel):
    question: str
    location: dict | None = None

# @app.post("/answer")
# async def run_query(body: QuestionRequest):
#     if not body.question.strip():
#         raise HTTPException(status_code=400, detail="Question is required")
#     response = await asyncio.to_thread(rag_chain_default.invoke, {"input": body.question})
#     return {"answer": response["answer"]}

@app.post("/answer")
async def run_query(body: QuestionRequest):
    if not body.question.strip():
        raise HTTPException(status_code=400, detail="Question is required")
    
    loc_str = ""
    if body.location:
        loc_str = f"User location: lat {body.location.get('lat')}, lon {body.location.get('lon')}"
        print(loc_str)

    response = await asyncio.to_thread(
        rag_chain_default.invoke,
        {
            "input": body.question,
            "location_info": loc_str
        }
    )
    return {"answer": response["answer"]}


>>>>>>> 809cc7dbcad3a17c96f6084921cee17d01c24b55
