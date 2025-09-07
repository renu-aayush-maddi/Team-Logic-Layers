# extractors.py
# Helpers for extracting text/chunks from multiple file types

import io
import os
import pandas as pd
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import logging

# OCR fallbacks
try:
    from google.cloud import vision
    GOOGLE_VISION_AVAILABLE = True
except Exception:
    GOOGLE_VISION_AVAILABLE = False

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

logger = logging.getLogger("extractors")

# ------------------------
# Utilities
# ------------------------
def chunk_text(text: str, max_chars: int = 2000):
    parts = []
    cur, cur_len = [], 0
    for para in [p.strip() for p in text.split("\n\n") if p.strip()]:
        if cur_len + len(para) + 2 > max_chars and cur:
            parts.append("\n\n".join(cur))
            cur, cur_len = [para], len(para)
        else:
            cur.append(para)
            cur_len += len(para) + 2
    if cur:
        parts.append("\n\n".join(cur))
    if not parts and text.strip():
        parts = [text.strip()[:max_chars]]
    return parts

# ------------------------
# PDF extraction (text + tables)
# ------------------------
def extract_pages_from_pdf_bytes_extended(pdf_bytes: bytes, filename: str = "uploaded.pdf"):
    out_chunks = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as e:
        logger.exception("Failed to open PDF: %s", e)
        return []

    for i in range(doc.page_count):
        page = doc.load_page(i)
        text = page.get_text("text") or ""
        for idx, c in enumerate(chunk_text(text)):
            out_chunks.append({
                "page": i + 1,
                "text": c,
                "filename": filename,
                "chunk_id": f"p{i+1}_c{idx+1}",
                "meta": {"source": "pdf"}
            })

    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for i, p in enumerate(pdf.pages):
                tables = p.extract_tables()
                for t_idx, table in enumerate(tables):
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0]) if len(table) > 1 else pd.DataFrame(table)
                        table_text = df.to_csv(index=False)
                        out_chunks.append({
                            "page": i + 1,
                            "text": f"[TABLE]\n" + table_text,
                            "filename": filename,
                            "chunk_id": f"p{i+1}_table{t_idx+1}",
                            "meta": {"source": "pdf_table"}
                        })
                    except Exception:
                        continue
    except Exception as e:
        logger.debug("pdfplumber failed: %s", e)

    return out_chunks

# ------------------------
# DOCX extraction
# ------------------------
def extract_text_from_docx_bytes(docx_bytes: bytes, filename: str = "doc.docx"):
    try:
        doc = Document(io.BytesIO(docx_bytes))
    except Exception as e:
        logger.exception("Failed to open docx: %s", e)
        return []
    full = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    chunks = chunk_text("\n\n".join(full))
    return [{"page": i+1, "text": c, "filename": filename, "chunk_id": f"docx_p{i+1}", "meta": {"source":"docx"}} for i, c in enumerate(chunks)]

# ------------------------
# PPTX extraction
# ------------------------
def extract_text_from_pptx_bytes(pptx_bytes: bytes, filename: str = "slides.pptx"):
    out = []
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        logger.exception("Failed to open pptx: %s", e)
        return []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shp in slide.shapes:
            try:
                if hasattr(shp, 'text') and shp.text:
                    parts.append(shp.text.strip())
            except Exception:
                continue
        for ci, c in enumerate(chunk_text("\n\n".join(parts))):
            out.append({
                "page": i+1,
                "text": c,
                "filename": filename,
                "chunk_id": f"slide{i+1}_c{ci+1}",
                "meta": {"source":"pptx"}
            })
    return out

# ------------------------
# Image OCR
# ------------------------
async def image_bytes_to_text(img_bytes: bytes, filename: str = "image.png", use_google: bool = False):
    text = ""
    if use_google and GOOGLE_VISION_AVAILABLE:
        try:
            client = vision.ImageAnnotatorClient()
            image = vision.Image(content=img_bytes)
            resp = client.document_text_detection(image=image)
            text = resp.full_text_annotation.text or ""
        except Exception as e:
            logger.exception("Google Vision OCR failed: %s", e)
            text = ""
    if not text and TESSERACT_AVAILABLE:
        try:
            im = Image.open(io.BytesIO(img_bytes)).convert('RGB')
            text = pytesseract.image_to_string(im)
        except Exception as e:
            logger.exception("Pytesseract OCR failed: %s", e)
            text = ""
    if not text:
        logger.warning("No OCR available for %s", filename)
        return []
    return [{"page": 1, "text": c, "filename": filename, "chunk_id": f"img_c{idx+1}", "meta": {"source":"image"}} for idx, c in enumerate(chunk_text(text))]

# ------------------------
# TXT extraction
# ------------------------
def extract_text_from_txt_bytes(txt_bytes: bytes, filename: str = "file.txt"):
    text = txt_bytes.decode('utf-8', errors='ignore')
    return [{"page": i+1, "text": c, "filename": filename, "chunk_id": f"txt_p{i+1}", "meta": {"source":"txt"}} for i, c in enumerate(chunk_text(text))]

# ------------------------
# CSV/Excel extraction
# ------------------------
def extract_tables_from_csv_bytes(csv_bytes: bytes, filename: str = "table.csv"):
    try:
        df = pd.read_csv(io.BytesIO(csv_bytes))
        text = df.to_csv(index=False)
        return [{"page": 1, "text": f"[TABLE]\n" + text, "filename": filename, "chunk_id": "csv_table_1", "meta": {"source": "csv"}}]
    except Exception as e:
        logger.exception("CSV parse failed: %s", e)
        return []

# ------------------------
# Unified ingest function (EXPORTED)
# ------------------------
async def ingest_file(upload_path: str, original_filename: str, use_google_vision: bool = False):
    """Main entry point: call this from main.py to get chunks for any supported file."""
    _, ext = os.path.splitext(original_filename.lower())
    with open(upload_path, 'rb') as f:
        b = f.read()

    if ext == '.pdf':
        return extract_pages_from_pdf_bytes_extended(b, filename=original_filename)
    if ext in ('.docx',):
        return extract_text_from_docx_bytes(b, filename=original_filename)
    if ext in ('.ppt', '.pptx'):
        return extract_text_from_pptx_bytes(b, filename=original_filename)
    if ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp'):
        return await image_bytes_to_text(b, filename=original_filename, use_google=use_google_vision)
    if ext == '.txt':
        return extract_text_from_txt_bytes(b, filename=original_filename)
    if ext == '.csv':
        return extract_tables_from_csv_bytes(b, filename=original_filename)
    if ext in ('.xls', '.xlsx'):
        try:
            df = pd.read_excel(io.BytesIO(b))
            text = df.to_csv(index=False)
            return [{"page": 1, "text": f"[TABLE]\n" + text, "filename": original_filename, "chunk_id": "excel_table_1", "meta": {"source": "excel"}}]
        except Exception as e:
            logger.exception("Excel parse failed: %s", e)
            return []

    logger.warning("Unsupported extension: %s", ext)
    return []

# For import in main.py
__all__ = ["ingest_file"]




# """
# extractors.py
# Unified file ingestion helpers for RAG pipeline.

# Supported:
# - PDF (PyMuPDF + pdfplumber table extraction)
# - DOCX
# - PPTX
# - Images (Google Vision optional, pytesseract fallback)
# - TXT
# - CSV / Excel (pandas)
# - Produces list[dict] chunks: { page, text, filename, chunk_id, meta }

# Usage:
# from extractors import ingest_file, chunk_text
# chunks = await ingest_file(saved_filepath, original_filename)
# """
# import os
# import io
# import logging
# from typing import List
# import uuid

# # text / doc libs
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation
# import pandas as pd
# import pdfplumber
# from PIL import Image

# # OCR fallbacks (optional)
# try:
#     from google.cloud import vision
#     GOOGLE_VISION_AVAILABLE = True
# except Exception:
#     GOOGLE_VISION_AVAILABLE = False

# try:
#     import pytesseract
#     TESSERACT_AVAILABLE = True
# except Exception:
#     TESSERACT_AVAILABLE = False

# logger = logging.getLogger("extractors")
# logger.setLevel(logging.DEBUG)

# # toggle via env var if desired
# USE_GOOGLE_VISION = os.getenv("USE_GOOGLE_VISION", "0") in ("1", "true", "True")


# def chunk_text(text: str, max_chars: int = 2000) -> List[str]:
#     """Naive paragraph chunker that keeps pieces <= max_chars."""
#     parts = []
#     cur = []
#     cur_len = 0
#     for para in [p.strip() for p in text.split("\n\n") if p.strip()]:
#         if cur_len + len(para) + 2 > max_chars and cur:
#             parts.append("\n\n".join(cur))
#             cur = [para]
#             cur_len = len(para)
#         else:
#             cur.append(para)
#             cur_len += len(para) + 2
#     if cur:
#         parts.append("\n\n".join(cur))
#     if not parts and text.strip():
#         parts = [text.strip()[:max_chars]]
#     return parts


# def extract_pages_from_pdf_bytes_extended(pdf_bytes: bytes, filename: str = "uploaded.pdf") -> List[dict]:
#     """Extract text chunks and attempt table extraction from a PDF."""
#     out_chunks = []
#     try:
#         doc = fitz.open(stream=pdf_bytes, filetype="pdf")
#     except Exception as e:
#         logger.exception("Failed to open PDF: %s", e)
#         return []

#     for i in range(doc.page_count):
#         page = doc.load_page(i)
#         text = page.get_text("text") or ""
#         chunks = chunk_text(text)
#         for idx, c in enumerate(chunks):
#             out_chunks.append({
#                 "page": i + 1,
#                 "text": c,
#                 "filename": filename,
#                 "chunk_id": f"{uuid.uuid4().hex}_p{i+1}_c{idx+1}",
#                 "meta": {"source": "pdf"}
#             })

#     # Attempt table extraction via pdfplumber
#     try:
#         with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
#             for i, p in enumerate(pdf.pages):
#                 tables = p.extract_tables()
#                 for t_idx, table in enumerate(tables):
#                     try:
#                         df = pd.DataFrame(table[1:], columns=table[0]) if len(table) > 1 else pd.DataFrame(table)
#                         table_text = df.to_csv(index=False)
#                         out_chunks.append({
#                             "page": i + 1,
#                             "text": "[TABLE]\n" + table_text,
#                             "filename": filename,
#                             "chunk_id": f"{uuid.uuid4().hex}_p{i+1}_table{t_idx+1}",
#                             "meta": {"source": "pdf_table"}
#                         })
#                     except Exception:
#                         continue
#     except Exception as e:
#         logger.debug("pdfplumber table extraction not available or failed: %s", e)

#     logger.info("PDF extraction produced %d chunks", len(out_chunks))
#     return out_chunks


# def extract_text_from_docx_bytes(docx_bytes: bytes, filename: str = "doc.docx") -> List[dict]:
#     try:
#         doc = Document(io.BytesIO(docx_bytes))
#     except Exception as e:
#         logger.exception("Failed to open docx bytes: %s", e)
#         return []
#     full = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
#     text = "\n\n".join(full)
#     chunks = chunk_text(text)
#     return [{
#         "page": i + 1,
#         "text": c,
#         "filename": filename,
#         "chunk_id": f"{uuid.uuid4().hex}_docx_p{i+1}",
#         "meta": {"source": "docx"}
#     } for i, c in enumerate(chunks)]


# def extract_text_from_pptx_bytes(pptx_bytes: bytes, filename: str = "slides.pptx") -> List[dict]:
#     out = []
#     try:
#         prs = Presentation(io.BytesIO(pptx_bytes))
#     except Exception as e:
#         logger.exception("Failed to open pptx bytes: %s", e)
#         return []
#     for i, slide in enumerate(prs.slides):
#         parts = []
#         for shp in slide.shapes:
#             try:
#                 if hasattr(shp, "text") and shp.text:
#                     parts.append(shp.text.strip())
#             except Exception:
#                 continue
#         text = "\n\n".join(parts)
#         chunks = chunk_text(text)
#         for ci, c in enumerate(chunks):
#             out.append({
#                 "page": i + 1,
#                 "text": c,
#                 "filename": filename,
#                 "chunk_id": f"{uuid.uuid4().hex}_slide{i+1}_c{ci+1}",
#                 "meta": {"source": "pptx"}
#             })
#     return out


# async def image_bytes_to_text(img_bytes: bytes, filename: str = "image.png") -> List[dict]:
#     """Try Google Vision if configured; otherwise use pytesseract fallback."""
#     text = ""
#     if USE_GOOGLE_VISION and GOOGLE_VISION_AVAILABLE:
#         try:
#             client = vision.ImageAnnotatorClient()
#             image = vision.Image(content=img_bytes)
#             resp = client.document_text_detection(image=image)
#             text = resp.full_text_annotation.text or ""
#             logger.debug("Google Vision OCR extracted %d chars", len(text))
#         except Exception as e:
#             logger.exception("Google Vision OCR failed, falling back: %s", e)
#             text = ""

#     if not text and TESSERACT_AVAILABLE:
#         try:
#             im = Image.open(io.BytesIO(img_bytes)).convert("RGB")
#             text = pytesseract.image_to_string(im)
#             logger.debug("Pytesseract extracted %d chars", len(text))
#         except Exception as e:
#             logger.exception("Pytesseract OCR failed: %s", e)
#             text = ""

#     if not text:
#         logger.warning("No OCR available to extract text from image %s", filename)
#         return []

#     chunks = chunk_text(text)
#     return [{
#         "page": 1,
#         "text": c,
#         "filename": filename,
#         "chunk_id": f"{uuid.uuid4().hex}_img_c{idx+1}",
#         "meta": {"source": "image"}
#     } for idx, c in enumerate(chunks)]


# def extract_text_from_txt_bytes(txt_bytes: bytes, filename: str = "file.txt") -> List[dict]:
#     text = txt_bytes.decode("utf-8", errors="ignore")
#     chunks = chunk_text(text)
#     return [{
#         "page": i + 1,
#         "text": c,
#         "filename": filename,
#         "chunk_id": f"{uuid.uuid4().hex}_txt_p{i+1}",
#         "meta": {"source": "txt"}
#     } for i, c in enumerate(chunks)]


# def extract_tables_from_csv_bytes(csv_bytes: bytes, filename: str = "table.csv") -> List[dict]:
#     try:
#         df = pd.read_csv(io.BytesIO(csv_bytes))
#         text = df.to_csv(index=False)
#         return [{
#             "page": 1,
#             "text": "[TABLE]\n" + text,
#             "filename": filename,
#             "chunk_id": f"{uuid.uuid4().hex}_csv_table_1",
#             "meta": {"source": "csv"}
#         }]
#     except Exception as e:
#         logger.exception("Failed to parse CSV: %s", e)
#         return []


# async def ingest_file(upload_path: str, original_filename: str) -> List[dict]:
#     """Dispatch file to the right extractor based on extension."""
#     _, ext = os.path.splitext(original_filename.lower())
#     with open(upload_path, "rb") as f:
#         b = f.read()

#     if ext == ".pdf":
#         return extract_pages_from_pdf_bytes_extended(b, filename=original_filename)
#     if ext in (".docx",):
#         return extract_text_from_docx_bytes(b, filename=original_filename)
#     if ext in (".ppt", ".pptx"):
#         return extract_text_from_pptx_bytes(b, filename=original_filename)
#     if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
#         return await image_bytes_to_text(b, filename=original_filename)
#     if ext in (".txt",):
#         return extract_text_from_txt_bytes(b, filename=original_filename)
#     if ext in (".csv",):
#         return extract_tables_from_csv_bytes(b, filename=original_filename)
#     if ext in (".xls", ".xlsx"):
#         try:
#             df = pd.read_excel(io.BytesIO(b))
#             text = df.to_csv(index=False)
#             return [{
#                 "page": 1,
#                 "text": "[TABLE]\n" + text,
#                 "filename": original_filename,
#                 "chunk_id": f"{uuid.uuid4().hex}_excel_table_1",
#                 "meta": {"source": "excel"}
#             }]
#         except Exception as e:
#             logger.exception("Failed to parse excel: %s", e)
#             return []

#     logger.warning("Unsupported extension for ingestion: %s", ext)
#     return []
